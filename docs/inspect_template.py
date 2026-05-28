"""Inspect the BashGU presentation template structure."""
import sys
io_enc = sys.stdout.encoding
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

p = Presentation('template.pptx')

print(f"Slide size: {p.slide_width} x {p.slide_height} EMU")
print(f"  = {p.slide_width / 914400:.2f}\" x {p.slide_height / 914400:.2f}\"")
print(f"Slides: {len(p.slides)}")
print(f"Layouts: {len(p.slide_layouts)}")
for i, layout in enumerate(p.slide_layouts):
    print(f"  Layout {i}: name='{layout.name}', placeholders={len(layout.placeholders)}")

print("\n=== Slides ===")
for i, slide in enumerate(p.slides):
    print(f"\n--- Slide {i+1} (layout: {slide.slide_layout.name}) ---")
    for sh in slide.shapes:
        info = f"  Shape: type={sh.shape_type}, name='{sh.name}'"
        if sh.has_text_frame and sh.text_frame.text:
            text = sh.text_frame.text.replace('\n', ' | ')
            info += f", text='{text[:80]}'"
        if hasattr(sh, 'left') and sh.left is not None:
            info += f", pos=({sh.left/914400:.2f}\", {sh.top/914400:.2f}\"), size=({sh.width/914400:.2f}\"x{sh.height/914400:.2f}\")"
        print(info)
        if sh.is_placeholder:
            print(f"    placeholder type={sh.placeholder_format.type}, idx={sh.placeholder_format.idx}")
